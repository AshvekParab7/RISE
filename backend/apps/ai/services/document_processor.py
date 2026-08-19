from pathlib import Path
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

def extract_document(resource):
    if not resource.file: return ''
    suffix = Path(resource.file.name).suffix.lower()
    path = resource.file.path
    if suffix == '.pdf': return '\n'.join(page.extract_text() or '' for page in PdfReader(path).pages)
    if suffix == '.docx': return '\n'.join(paragraph.text for paragraph in Document(path).paragraphs)
    if suffix == '.pptx': return '\n'.join(shape.text for slide in Presentation(path).slides for shape in slide.shapes if hasattr(shape, 'text'))
    if suffix == '.txt': return Path(path).read_text(encoding='utf-8')
    raise ValueError('This file type cannot be processed yet.')
