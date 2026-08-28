from html import escape
from io import BytesIO
from pathlib import Path


class UnsupportedOfficePreview(Exception):
    pass


def _paragraph_html(paragraph):
    text = escape(paragraph.text.strip())
    if not text:
        return ''
    style = paragraph.style.name.lower() if paragraph.style else ''
    if 'heading' in style:
        level = next((value for value in ('1', '2', '3') if value in style), '3')
        return f'<h{level}>{text}</h{level}>'
    return f'<p>{text}</p>'


def _docx_html(content):
    from docx import Document

    document = Document(BytesIO(content))
    blocks = [_paragraph_html(paragraph) for paragraph in document.paragraphs]
    for table in document.tables:
        rows = []
        for row in table.rows:
            cells = ''.join(f'<td>{escape(cell.text.strip())}</td>' for cell in row.cells)
            rows.append(f'<tr>{cells}</tr>')
        blocks.append(f'<table><tbody>{"".join(rows)}</tbody></table>')
    return ''.join(block for block in blocks if block)


def _pptx_html(content):
    from pptx import Presentation

    presentation = Presentation(BytesIO(content))
    slides = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if not getattr(shape, 'has_text_frame', False):
                continue
            text = escape(shape.text.strip())
            if text:
                parts.append(f'<p>{text.replace(chr(10), "<br>")}</p>')
        slides.append(f'<section><h2>Slide {index}</h2>{"".join(parts)}</section>')
    return ''.join(slides)


def office_preview_html(content, filename, mime_type=''):
    suffix = Path(filename).suffix.lower()
    if suffix == '.docx' or mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
        body = _docx_html(content)
    elif suffix == '.pptx' or mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
        body = _pptx_html(content)
    else:
        raise UnsupportedOfficePreview
    return f'''<!doctype html>
<html><head><meta charset="utf-8"><style>
body {{ margin: 0; padding: 28px; color: #302d34; background: #fff; font: 16px/1.6 Georgia, serif; }}
main {{ max-width: 850px; margin: auto; }}
h1, h2, h3 {{ line-height: 1.2; }}
table {{ width: 100%; margin: 18px 0; border-collapse: collapse; }}
td {{ padding: 8px; border: 1px solid #d8d0da; vertical-align: top; }}
section {{ margin: 0 0 30px; padding-bottom: 20px; border-bottom: 1px solid #e5dfe7; }}
</style></head><body><main>{body or '<p>This document contains no readable text.</p>'}</main></body></html>'''
